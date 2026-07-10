import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_first_slide():
    prs = Presentation()

    # Use a blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Background color - White (Clean Tech Enterprise Theme)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.add_paragraph()
    p.text = "Next-Generation Web Ecosystem &\nAI Chatbot for ViBeS Lab"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.name = 'Segoe UI'
    p.font.color.rgb = RGBColor(10, 37, 64) # Deep Navy Blue

    # Subtitle
    left = Inches(1)
    top = Inches(4.2)
    width = Inches(8)
    height = Inches(1.0)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    p2 = tf2.add_paragraph()
    p2.text = "A decoupled, intelligent digital portal for modern academic research"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(24)
    p2.font.name = 'Segoe UI'
    p2.font.color.rgb = RGBColor(0, 212, 178) # Vibrant Teal

    # Footer
    left = Inches(5.5)
    top = Inches(6.5)
    width = Inches(4)
    height = Inches(0.5)
    txBox3 = slide.shapes.add_textbox(left, top, width, height)
    tf3 = txBox3.text_frame

    p3 = tf3.add_paragraph()
    p3.text = "Presented by: [Your Name] | July 2026"
    p3.alignment = PP_ALIGN.RIGHT
    p3.font.size = Pt(14)
    p3.font.name = 'Segoe UI'
    p3.font.color.rgb = RGBColor(128, 128, 128) # Medium gray

    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), 'ViBeS_Presentation_Slide1.pptx')
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == '__main__':
    create_first_slide()
